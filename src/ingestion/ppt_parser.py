"""PowerPoint ingestion: slide text, tables, speaker notes, image-only slide flags.

.ppt (the pre-2007 binary format) cannot be read by python-pptx. Rather than failing
silently we say so and, where LibreOffice is available, convert it first.
"""

from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path

from ..models.evidence import SourceType, Warning
from ..utils.config import IMAGE_PAGE_CHAR_THRESHOLD
from ..utils.logging import get_logger
from ..utils.text import clean
from .types import ParsedDocument, Segment, UnreadableFile

_log = get_logger()


def parse_pptx(path: str | Path, source_type: SourceType = SourceType.PITCH_DECK) -> ParsedDocument:
    path = Path(path)
    converted: Path | None = None
    warnings: list[Warning] = []

    if path.suffix.lower() == ".ppt":
        converted = _convert_ppt(path)
        if converted is None:
            raise UnreadableFile(
                f"{path.name} is a legacy .ppt file and no LibreOffice install was found to "
                "convert it. Re-save it as .pptx or .pdf and try again."
            )
        warnings.append(
            Warning(
                severity="info",
                stage="ingestion",
                message=f"{path.name} was converted from .ppt to .pptx before reading.",
            )
        )

    target = converted or path

    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - dependency is declared
        raise UnreadableFile("python-pptx is not installed.") from exc

    try:
        presentation = Presentation(str(target))
    except Exception as exc:
        raise UnreadableFile(f"{path.name} could not be opened as a PowerPoint file: {exc}") from exc

    doc = ParsedDocument(path=path, name=path.name, source_type=source_type, kind="deck")
    doc.warnings = warnings

    for number, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        tables: list[list[list[str]]] = []
        title = ""
        picture_count = 0

        for shape in slide.shapes:
            try:
                if shape.has_table:
                    rows = [[clean(cell.text) for cell in row.cells] for row in shape.table.rows]
                    if any(any(cell for cell in row) for row in rows):
                        tables.append(rows)
                    continue
            except Exception:
                pass
            if getattr(shape, "shape_type", None) is not None and shape.shape_type == 13:
                picture_count += 1
            try:
                if shape.has_text_frame:
                    text = clean(shape.text_frame.text)
                    if text:
                        texts.append(text)
            except Exception:
                continue

        try:
            if slide.shapes.title is not None:
                title = clean(slide.shapes.title.text)
        except Exception:
            title = ""

        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = clean(slide.notes_slide.notes_text_frame.text)
        except Exception:
            notes = ""

        segment = Segment(
            index=number,
            text="\n".join(texts),
            kind="slide",
            title=title,
            notes=notes,
            tables=tables,
        )
        segment.image_only = segment.char_count < IMAGE_PAGE_CHAR_THRESHOLD and picture_count > 0
        doc.segments.append(segment)

    image_slides = [s.index for s in doc.segments if s.image_only]
    if image_slides:
        preview = ", ".join(str(i) for i in image_slides[:8])
        more = "..." if len(image_slides) > 8 else ""
        doc.warnings.append(
            Warning(
                severity="warning",
                stage="ingestion",
                message=(
                    f"Possible image-only content on slide {preview}{more} of {path.name} - "
                    "manual verification recommended."
                ),
            )
        )
    doc.metadata = {"slides": len(doc.segments), "reader": "python-pptx"}
    return doc


def _convert_ppt(path: Path) -> Path | None:
    """Convert a legacy .ppt to .pptx with LibreOffice, if one is installed."""
    soffice = _soffice()
    if not soffice:
        return None
    outdir = Path(tempfile.mkdtemp(prefix="lim_ppt_"))
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pptx", "--outdir", str(outdir), str(path)],
            check=True,
            capture_output=True,
            timeout=180,
        )
    except Exception as exc:
        _log.debug("LibreOffice conversion failed for %s: %s", path.name, exc)
        return None
    candidates = list(outdir.glob("*.pptx"))
    return candidates[0] if candidates else None


def _soffice() -> str | None:
    for name in ("soffice", "soffice.exe", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    for candidate in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ):
        if Path(candidate).exists():
            return candidate
    return None
