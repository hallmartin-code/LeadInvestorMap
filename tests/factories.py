"""Builders for test fixtures, so each test states only what it is actually about."""

from __future__ import annotations

from pathlib import Path

from src.models.company import Company
from src.models.evidence import Confidence, Fact, SourceRef, SourceType
from src.models.investor import (
    DiligenceStage,
    FundStatus,
    Investor,
    InvestorType,
    LeadHistoryEntry,
    Relationship,
)
from src.models.round import Round


def source(name: str = "deck.pdf", page: int = 1, text: str = "") -> SourceRef:
    return SourceRef(
        source_type=SourceType.PITCH_DECK, source_name=name, page_or_slide=page, source_text=text
    )


def fact(value: str, numeric: float | None = None, claim: str = "claim") -> Fact:
    return Fact.from_document(claim, value, source(), numeric_value=numeric)


def make_round(
    stage: str = "Series A",
    raise_amount: float = 6_000_000,
    committed: float | None = 1_500_000,
    target_close: str | None = "December 2026",
) -> Round:
    round_ = Round(
        stage=fact(stage, claim="Round stage"),
        raise_amount=fact(f"${raise_amount / 1e6:g}M", raise_amount, "Round size"),
        instrument=fact("Priced Equity", claim="Instrument"),
    )
    if committed is not None:
        round_.committed = fact(f"${committed / 1e6:g}M", committed, "Amount committed")
    if target_close is not None:
        round_.target_close = fact(target_close, claim="Target close")
    return round_


def make_company(
    name: str = "Testco",
    sector: str = "Healthcare / Life Sciences",
    keywords: list[str] | None = None,
    competitors: list[str] | None = None,
) -> Company:
    return Company(
        name=fact(name, claim="Company name"),
        one_liner=fact("A test company", claim="Company one-liner"),
        sector=fact(sector, claim="Sector"),
        keywords=keywords or ["diagnostics", "clinical", "sepsis"],
        named_competitors=competitors or [],
    )


def make_investor(
    name: str,
    *,
    investor_type: InvestorType = InvestorType.VC,
    check: tuple[float | None, float | None] = (2_000_000, 5_000_000),
    lead_history: list[tuple[str, str, str]] | None = None,
    leads_rounds_stated: bool | None = None,
    entry_stages: list[str] | None = None,
    sector_focus: str = "diagnostics, clinical tools",
    fund_status: FundStatus = FundStatus.ACTIVE,
    relationship: Relationship = Relationship.FIRST_MEETING,
    diligence: DiligenceStage = DiligenceStage.FIRST_MEETING,
    portfolio: list[str] | None = None,
    warm_intro: str | None = None,
    notes: str = "",
    stated_dependencies: list[str] | None = None,
) -> Investor:
    investor = Investor(
        investor_name=name,
        investor_type=investor_type,
        estimated_check_min=check[0],
        estimated_check_max=check[1],
        leads_rounds_stated=leads_rounds_stated,
        entry_stages=entry_stages if entry_stages is not None else ["Series A"],
        sector_fit_detail=f"Stated focus: {sector_focus}" if sector_focus else "",
        fund_status=fund_status,
        relationship_strength=relationship,
        current_diligence_stage=diligence,
        supporting_portfolio_companies=portfolio or [],
        warm_intro_path=warm_intro,
        warm_intro_verified=bool(warm_intro),
        notes=notes,
        stated_dependencies=stated_dependencies or [],
        confidence=Confidence.MEDIUM,
    )
    investor.add_source(
        SourceRef(
            source_type=SourceType.INVESTOR_LIST,
            source_name="targets.csv",
            page_or_slide=2,
            source_text=name,
        )
    )
    for company, round_label, role in lead_history or []:
        investor.lead_history.append(
            LeadHistoryEntry(
                company=company,
                round_label=round_label,
                role=role,
                year="2025",
                source=investor.sources[0],
                confidence=Confidence.MEDIUM,
            )
        )
    return investor


def write_csv(path: Path, rows: list[dict]) -> Path:
    import csv

    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows(rows)
    return path


def write_pdf(path: Path, pages: list[list[str]]) -> Path:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas as pdfcanvas

    canvas = pdfcanvas.Canvas(str(path), pagesize=letter)
    _, height = letter
    for lines in pages:
        y = height - 80
        canvas.setFont("Helvetica", 12)
        for line in lines:
            canvas.drawString(56, y, line)
            y -= 18
        canvas.showPage()
    canvas.save()
    return path


def write_pptx(path: Path, slides: list[tuple[str, list[str], str]]) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    blank = presentation.slide_layouts[5]
    for title, lines, notes in slides:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
        frame = box.text_frame
        frame.text = lines[0] if lines else ""
        for line in lines[1:]:
            frame.add_paragraph().text = line
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    presentation.save(str(path))
    return path
